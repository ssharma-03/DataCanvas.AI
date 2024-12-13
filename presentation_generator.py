presentation_generation.py 


# presentation_generator.py
import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
import seaborn as sns
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import io
import tempfile

class PresentationGenerator:
    def __init__(self, style=None):
        self.style = style or {
            'background': '#ffffff',
            'text': '#000000',
            'accent': '#2c3e50',
            'font': 'arial'
        }
        
    def update_style(self, style):
        """Update the presentation style"""
        self.style = style
        
    def create_presentation(self, data, columns, selected_charts, title="Data Analysis Report", 
                          company_name="", include_stats=True, include_conclusions=True):
        """Create PowerPoint presentation with data analysis and charts"""
        try:
            prs = Presentation()
            
            # Title slide
            title_slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(title_slide_layout)
            
            # Set title
            title_shape = slide.shapes.title
            title_shape.text = title
            
            # Set subtitle if company name is provided
            if company_name and len(slide.placeholders) > 1:
                subtitle_shape = slide.placeholders[1]
                subtitle_shape.text = company_name
            
            # Add content slides
            self._add_overview_slide(prs, data, columns, selected_charts)
            
            if include_stats:
                for col in columns:
                    self._add_stats_slide(prs, data, col)
            
            for chart_type in selected_charts:
                self._add_chart_slide(prs, data, columns, chart_type)
            
            if include_conclusions:
                self._add_conclusions_slide(prs, data, columns)
            
            # Save to BytesIO
            output = io.BytesIO()
            prs.save(output)
            output.seek(0)
            
            return output.getvalue()
            
        except Exception as e:
            raise Exception(f"Error creating presentation: {str(e)}")

    def _add_overview_slide(self, prs, data, columns, selected_charts):
        """Add overview slide to presentation"""
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # Using blank layout
        
        # Add title
        title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = "Data Overview"
        
        # Add content
        content_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5))
        text_frame = content_box.text_frame
        
        p = text_frame.add_paragraph()
        p.text = "Dataset Information:"
        
        p = text_frame.add_paragraph()
        p.text = f"• Number of Variables: {len(columns)}"
        
        p = text_frame.add_paragraph()
        p.text = f"• Total Records: {len(data)}"
        
        p = text_frame.add_paragraph()
        p.text = f"• Time Period: {len(data)} points"
        
        p = text_frame.add_paragraph()
        p.text = "\nVariables Analyzed:"
        
        for col in columns:
            p = text_frame.add_paragraph()
            p.text = f"• {col}"
        
        p = text_frame.add_paragraph()
        p.text = "\nVisualizations Created:"
        
        for chart in selected_charts:
            p = text_frame.add_paragraph()
            p.text = f"• {chart}"

    def _add_stats_slide(self, prs, data, column):
        """Add statistical analysis slide for a column"""
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # Using blank layout
        
        # Add title
        title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = f"{column} - Statistical Analysis"
        
        # Add content
        content_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5))
        text_frame = content_box.text_frame
        
        stats = data[column].describe()
        change = ((data[column].iloc[-1] / data[column].iloc[0] - 1) * 100)
        
        p = text_frame.add_paragraph()
        p.text = "Basic Statistics:"
        
        stats_items = [
            f"• Mean: {stats['mean']:.2f}",
            f"• Median: {data[column].median():.2f}",
            f"• Standard Deviation: {stats['std']:.2f}",
            f"• Minimum: {stats['min']:.2f}",
            f"• Maximum: {stats['max']:.2f}",
            "\nDistribution Analysis:",
            f"• 25th Percentile: {stats['25%']:.2f}",
            f"• 75th Percentile: {stats['75%']:.2f}",
            f"• IQR: {(stats['75%'] - stats['25%']):.2f}",
            "\nChange Analysis:",
            f"• Total Change: {(data[column].iloc[-1] - data[column].iloc[0]):.2f}",
            f"• Percentage Change: {change:.2f}%"
        ]
        
        for item in stats_items:
            p = text_frame.add_paragraph()
            p.text = item

    def _add_chart_slide(self, prs, data, columns, chart_type):
        """Add visualization slide"""
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # Using blank layout
        
        # Add title
        title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = f"{chart_type} Analysis"
        
        # Create chart
        plt.figure(figsize=(10, 6))
        plt.style.use('dark_background' if self.style['background'] == '#000000' else 'default')
        
        self._create_chart(data, columns, chart_type)
        
        # Save chart to bytes
        img_stream = io.BytesIO()
        plt.savefig(img_stream, format='png', 
                   facecolor=self.style['background'],
                   bbox_inches='tight',
                   dpi=300)
        plt.close()
        
        # Add chart to slide
        img_stream.seek(0)
        slide.shapes.add_picture(img_stream, Inches(1), Inches(2), width=Inches(8))

    def _create_chart(self, data, columns, chart_type):
        """Create specific chart type"""
        if chart_type == 'Line Plot':
            for col in columns:
                plt.plot(range(len(data)), data[col], label=col, linewidth=2)
        elif chart_type == 'Bar Chart':
            data[columns].plot(kind='bar')
        elif chart_type == 'Scatter Plot':
            for col in columns:
                plt.scatter(range(len(data)), data[col], label=col, alpha=0.7)
        elif chart_type == 'Box Plot':
            plt.boxplot([data[col] for col in columns], labels=columns)
        elif chart_type == 'Violin Plot':
            plt.violinplot([data[col] for col in columns])
            plt.xticks(range(1, len(columns) + 1), columns)
        elif chart_type == 'Heatmap':
            sns.heatmap(data[columns].corr(), annot=True, cmap='coolwarm')
            
        plt.title(chart_type)
        plt.xlabel("Time Period")
        plt.ylabel("Value")
        if chart_type not in ['Heatmap', 'Box Plot', 'Violin Plot']:
            plt.legend()
        plt.grid(True, alpha=0.3)

    def _add_conclusions_slide(self, prs, data, columns):
        """Add conclusions slide"""
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # Using blank layout
        
        # Add title
        title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = "Key Findings"
        
        # Add content
        content_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5))
        text_frame = content_box.text_frame
        
        p = text_frame.add_paragraph()
        p.text = "Summary:"
        
        p = text_frame.add_paragraph()
        p.text = f"• Analyzed {len(columns)} variables over {len(data)} time points"
        
        p = text_frame.add_paragraph()
        p.text = "\nKey Metrics:"
        
        for col in columns:
            p = text_frame.add_paragraph()
            p.text = (f"• {col}: Mean = {data[col].mean():.2f}, "
                     f"Change = {((data[col].iloc[-1] / data[col].iloc[0] - 1) * 100):+.2f}%")
